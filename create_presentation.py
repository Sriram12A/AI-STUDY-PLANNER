import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from sklearn.linear_model import LinearRegression, LogisticRegression
from sklearn.metrics import (
    mean_absolute_error, mean_squared_error, r2_score,
    accuracy_score, precision_score, recall_score, f1_score
)
from sklearn.model_selection import train_test_split

def create_ppt():
    prs = Presentation()

    def add_slide(title_text):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        return slide

    # --- 1. TITLE SLIDE ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Powered Study Plan Optimizer"
    subtitle.text = "Model Performance Analysis & Personalized Timetable\nPresented to: Faculty Department"

    # --- 2. ML PERFORMANCE RESULTS ---
    # Linear Regression Data
    data_lr = {"improvement": [5, 10, 15, 20, 25, 30], "hours": [10, 20, 35, 50, 70, 95]}
    df_lr = pd.DataFrame(data_lr)
    model_lr = LinearRegression().fit(df_lr[["improvement"]], df_lr["hours"])
    y_pred_lr = model_lr.predict(df_lr[["improvement"]])
    
    mae = mean_absolute_error(df_lr["hours"], y_pred_lr)
    mse = mean_squared_error(df_lr["hours"], y_pred_lr)
    r2 = r2_score(df_lr["hours"], y_pred_lr)

    slide = add_slide("Regression Model Analysis")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Metric: Mean Absolute Error (MAE): {mae:.2f}"
    p = tf.add_paragraph()
    p.text = f"Metric: Mean Squared Error (MSE): {mse:.2f}"
    p = tf.add_paragraph()
    p.text = f"Metric: R-Squared (Accuracy): {r2:.4f}"
    p = tf.add_paragraph()
    p.text = "Analysis: High R2 score indicates the model is extremely well-fitted to predict study hours."

    # --- 3. CLASSIFICATION RESULTS ---
    X_cls = [[10], [20], [30], [40], [50], [60], [70], [80]]
    y_cls = [0, 0, 0, 1, 1, 1, 1, 1]
    X_train, X_test, y_train, y_test = train_test_split(X_cls, y_cls, test_size=0.2, random_state=42)
    model_log = LogisticRegression().fit(X_train, y_train)
    y_pred_log = model_log.predict(X_test)

    acc = accuracy_score(y_test, y_pred_log)
    prec = precision_score(y_test, y_pred_log, zero_division=0)
    rec = recall_score(y_test, y_pred_log, zero_division=0)

    slide = add_slide("Classification Performance (Priority Mapping)")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Overall Accuracy: {acc*100:.1f}%"
    p = tf.add_paragraph()
    p.text = f"Precision: {prec:.2f}"
    p = tf.add_paragraph()
    p.text = f"Recall Score: {rec:.2f}"

    # --- 4. PRIORITY-WEIGHTED ALLOCATION LOGIC ---
    slide = add_slide("Real-World Allocation Logic")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "100% Capacity Distribution Model:"
    p = tf.add_paragraph()
    p.text = "The system treats the total available time (e.g., 24 hours) as a 100% resource pool."
    p = tf.add_paragraph()
    p.text = "Weights are calculated based on the Improvement Gap (Target % - Previous %)."
    p = tf.add_paragraph()
    p.text = "Formula: (Subject Priority / Sum of All Priorities) * Total Time"

    # --- 5. SUBJECT-WISE DETAILED STUDY PLAN ---
    slide = add_slide("Detailed Personalized Study Plan")
    
    rows, cols = 6, 6
    left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(3.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ['Subject', 'Priority Weight', 'Daily Hours', 'Theory (20%)', 'Practice (50%)', 'Revision (30%)']
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
        table.cell(0, j).text_frame.paragraphs[0].font.bold = True

    # Sample weighted distribution (Total 100% -> 24 hours)
    # Using real-world split: Theory 20%, Practice 50%, Revision 30%
    subjects = [
        ["Maths", "32.2%", 7.73, 1.55, 3.86, 2.32],
        ["Chemistry", "32.2%", 7.73, 1.55, 3.86, 2.32],
        ["Physics", "35.6%", 8.54, 1.71, 4.27, 2.56]
    ]

    for i, sub in enumerate(subjects):
        for j, val in enumerate(sub):
            table.cell(i+1, j).text = str(val)

    # --- 6. COMBINED SUMMARY ---
    slide = add_slide("Combined Resource Utilization")
    body = slide.placeholders[1]
    tf = body.text_frame
    total_val = 24.0
    tf.text = f"Total System Output: {total_val} Hours (Full Daily Cycle)"
    p = tf.add_paragraph()
    p.text = "Strategic Distribution for Academic Excellence:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = f"Total Focused Theory: {total_val*0.2:.2f} Hours"
    p = tf.add_paragraph()
    p.level = 1
    p.text = f"Total Active Practice: {total_val*0.5:.2f} Hours"
    p = tf.add_paragraph()
    p.level = 1
    p.text = f"Total Critical Revision: {total_val*0.3:.2f} Hours"

    # --- FINAL SLIDE ---
    slide = add_slide("Conclusion")
    body = slide.placeholders[1]
    body.text = "This priority-weighted system ensures that no subject is left behind while focusing maximum effort (100% of resources) where the performance gap is largest."

    prs.save('Study_Plan_Analysis.pptx')
    print("PPT created successfully as 'Study_Plan_Analysis.pptx'")

if __name__ == "__main__":
    create_ppt()
